import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field, field_validator

from src.constants import PPTX_STYLING, SYSTEM_PROMPTS
from src.services.gemini_service import GeminiService
from src.utils.logger import logger


class SlideContent(BaseModel):
    title: str = Field(..., description="Slide title")
    bullets: list[str] = Field(
        ..., description="Max 4-5 concise bullet points (each under 100 characters)"
    )
    speaker_notes: str = Field(..., description="Detailed speaker notes explaining these points")

    @field_validator("bullets")
    @classmethod
    def validate_and_truncate_bullets(cls, v: list[str]) -> list[str]:
        # Limit to max 5 bullets per slide
        v = v[:5]
        cleaned_bullets = []
        for b in v:
            # Strip duplicate bullet characters (-, *, •)
            cleaned = b.lstrip("-*• ").strip()
            # Truncate long bullets (max 120 chars) to prevent layout overflow
            if len(cleaned) > 120:
                cleaned = cleaned[:117] + "..."
            cleaned_bullets.append(cleaned)
        return cleaned_bullets


class PresentationContent(BaseModel):
    title: str = Field(..., description="Main presentation title")
    subtitle: str = Field(..., description="Subtitle detailing target audience and date")
    agenda: list[str] = Field(..., description="List of agenda items")
    slides: list[SlideContent] = Field(
        ...,
        description="List of content slides (Exec Summary, Findings, Analysis, Recs, Conclusion, References)",
    )


class SlideService:
    """Service to generate slide deck structure and render PPTX presentations."""

    def __init__(self, gemini_service: GeminiService):
        self.gemini_service = gemini_service

    def generate_presentation(
        self, report_markdown: str, slide_count: int = 8, model: str | None = None
    ) -> PresentationContent:
        """Uses Gemini to summarize the final reviewed report into structured slide content."""
        logger.info("Summarizing report into presentation slide content...")

        prompt = (
            f"Please convert the following report into a presentation slide deck with approximately {slide_count} slides.\n\n"
            f"--- REPORT CONTENT ---\n{report_markdown}\n\n"
            "Generate slide content following the structure of the report, including Title, Agenda, Executive Summary, "
            "Key Findings, Analysis, Recommendations, Conclusion, and References. "
            "Keep slide bullet points extremely concise to prevent layout overflow. "
            "Provide detailed speaker notes for every slide."
        )

        system_instruction = SYSTEM_PROMPTS["slide_generator"]

        # Call structured generation
        presentation_content = self.gemini_service.generate_structured(
            prompt=prompt,
            response_schema=PresentationContent,
            system_instruction=system_instruction,
            model=model,
            temperature=0.3,
        )

        return presentation_content

    def build_pptx(self, content: PresentationContent, output_path: str) -> None:
        """Generates a professional 16:9 PPTX file using python-pptx from structured content."""
        logger.info(f"Building PPTX presentation file at: {output_path}")

        prs = Presentation()
        # Set widescreen slide size (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        # Helper to set slide background
        def set_bg(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*PPTX_STYLING["bg_color"])

        # 1. TITLE SLIDE
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide)

        # Main Title Box (centered)
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
        tf = title_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = content.title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = PPTX_STYLING["font_title"]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*PPTX_STYLING["primary_color"])

        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = content.subtitle
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = PPTX_STYLING["font_body"]
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(*PPTX_STYLING["secondary_color"])
        p2.space_before = Pt(20)

        # Set speaker notes dynamically
        slide.notes_slide.notes_text_frame.text = (
            f"Presentation: {content.title}\n"
            f"Context / Subtitle: {content.subtitle}\n"
            f"Introduction: Welcome the audience and outline the goals of today's review."
        )

        # 2. AGENDA SLIDE
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide)

        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        t_tf = t_box.text_frame
        t_p = t_tf.paragraphs[0]
        t_p.text = "Agenda"
        t_p.font.name = PPTX_STYLING["font_title"]
        t_p.font.size = Pt(36)
        t_p.font.bold = True
        t_p.font.color.rgb = RGBColor(*PPTX_STYLING["primary_color"])

        # Agenda bullets
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.0), Inches(10.333), Inches(4.5)
        )
        c_tf = content_box.text_frame
        c_tf.word_wrap = True

        clean_agenda_items = []
        for i, item in enumerate(content.agenda):
            # Clean duplicate numbering prefix
            clean_item = re.sub(r"^\d+[\.\s\-]+", "", item).strip()
            clean_agenda_items.append(clean_item)
            p = c_tf.paragraphs[0] if i == 0 else c_tf.add_paragraph()
            p.text = f"{i + 1}. {clean_item}"
            p.font.name = PPTX_STYLING["font_body"]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*PPTX_STYLING["text_color"])
            p.space_after = Pt(14)

        agenda_str = ", ".join(clean_agenda_items)
        slide.notes_slide.notes_text_frame.text = f"Agenda Overview:\nWe will cover the following sections in today's presentation: {agenda_str}."

        # 3. CONTENT SLIDES
        for slide_data in content.slides:
            slide = prs.slides.add_slide(blank_layout)
            set_bg(slide)

            # Slide Title
            t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
            t_tf = t_box.text_frame
            t_p = t_tf.paragraphs[0]
            t_p.text = slide_data.title
            t_p.font.name = PPTX_STYLING["font_title"]
            t_p.font.size = Pt(32)
            t_p.font.bold = True
            t_p.font.color.rgb = RGBColor(*PPTX_STYLING["primary_color"])

            # Bullets
            content_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(2.0), Inches(10.933), Inches(4.5)
            )
            c_tf = content_box.text_frame
            c_tf.word_wrap = True

            for i, bullet in enumerate(slide_data.bullets):
                # Bullets are already clean due to SlideContent validation
                p = c_tf.paragraphs[0] if i == 0 else c_tf.add_paragraph()
                p.text = f"•  {bullet}"
                p.font.name = PPTX_STYLING["font_body"]
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(*PPTX_STYLING["text_color"])
                p.space_after = Pt(12)
                p.level = 0

            # Set speaker notes
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

        prs.save(output_path)
        logger.info(f"PPTX successfully written to: {output_path}")

    def generate_speaker_notes_text(self, content: PresentationContent) -> str:
        """Formats the speaker notes into a clean, downloadable markdown document."""
        notes = []
        notes.append(f"# Speaker Notes: {content.title}")
        notes.append(f"Subtitle: {content.subtitle}\n")
        notes.append("## Agenda Slide")

        agenda_items = [re.sub(r"^\d+[\.\s\-]+", "", item).strip() for item in content.agenda]
        agenda_str = ", ".join(agenda_items)
        notes.append(f"We will guide you through the following topics: {agenda_str}.\n")

        for slide in content.slides:
            notes.append(f"## Slide: {slide.title}")
            notes.append(f"{slide.speaker_notes}\n")

        return "\n".join(notes)
